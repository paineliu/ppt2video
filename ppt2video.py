from pptx import Presentation
from pptx_tools import utils
import os
import edge_tts
import asyncio
import cv2
import shutil
import pandas as pd
from mutagen.mp3 import MP3
from moviepy.editor import VideoFileClip, AudioFileClip, concatenate_videoclips

def text_to_mp3(text, mp3_filename):
    print(mp3_filename)
    if os.path.isfile(mp3_filename):
        return
    
    voice = 'zh-CN-XiaoxiaoNeural'
    rate = '-4%'
    volume = '+0%'

    async def tts(text, mp3_filename):
        tts = edge_tts.Communicate(text=text, voice=voice, rate=rate, volume=volume)
        await tts.save(mp3_filename)

    asyncio.run(tts(text, mp3_filename))

def img_to_silent(img_filename, silent_filename, duration):
    print(silent_filename, duration)
    os.makedirs(os.path.dirname(silent_filename), exist_ok=True)
    
    fps = 1      # 帧率
    img_size = (1920, 1080)      # 图片尺寸
    fourcc = cv2.VideoWriter_fourcc(*"mp4v")
    videoWriter = cv2.VideoWriter(silent_filename, fourcc, fps, img_size)
    for i in range(0, duration):
        frame = cv2.imread(img_filename)
        frame = cv2.resize(frame, img_size)  # 生成视频   图片尺寸和设定尺寸相同
        videoWriter.write(frame)  # 将图片写进视频里
    videoWriter.release()  # 释放资源

def get_mp3_duration(mp3_filename):
    audio = MP3(mp3_filename)
    time_count = int(audio.info.length)
    return time_count

def slient_audio_to_video(silent_filename, mp3_filename, mp4_filename):
    print(mp4_filename)
    os.makedirs(os.path.dirname(mp4_filename), exist_ok=True)
    video = VideoFileClip(silent_filename)
    videos = video.set_audio(AudioFileClip(mp3_filename)) 
    videos.write_videofile(mp4_filename, audio_codec='aac') 

def merge_vidio(video_files, output_filename):
    print(output_filename)
    os.makedirs(os.path.dirname(output_filename), exist_ok=True)
    clip_lst = []
    for video in video_files:
        clip_lst.append(VideoFileClip(video))
    final_clip = concatenate_videoclips(clip_lst)
    final_clip.write_videofile(output_filename, codec="libx264", fps=24, audio_bitrate="50k") 

def cvs_to_mp3(cvs_filename, mp3_pathname):
    os.makedirs(mp3_pathname, exist_ok=True)
    out_data = pd.read_csv(cvs_filename)
    for i in range(len(out_data)):
        text  = out_data['text'][i]
        mp3_filename = os.path.join(mp3_pathname, '{}.mp3'.format(i + 1))
        text_to_mp3(text, mp3_filename)
    return len(out_data)

def make_video(res_pathname, temp_pathname, output_filename):

    cvs_filename = os.path.join(res_pathname, 'text.csv')
    mp3_pathname = os.path.join(temp_pathname, 'mp3')
    mp3_total = cvs_to_mp3(cvs_filename, mp3_pathname)

    img_pathname = os.path.join(temp_pathname, 'image')
    ppt_filename = os.path.join(res_pathname, 'slide.pptx')
    slide_to_image(ppt_filename, img_pathname)

    silent_pathname = os.path.join(temp_pathname, 'silent')
    
    video_pathname = os.path.join(temp_pathname, 'video')
    video_list = []
    for i in range(1, mp3_total + 1):
        mp3_filename = os.path.join(mp3_pathname, '{}.mp3'.format(i))
        mp3_duration = get_mp3_duration(mp3_filename)
        silent_filename = os.path.join(silent_pathname, '{}.mp4'.format(i))
        img_zh_filename = os.path.join(img_pathname, '幻灯片{}.PNG'.format(i))
        img_filename = os.path.join(img_pathname, '{}.PNG'.format(i))
        os.rename(img_zh_filename, img_filename)
        img_filename = os.path.abspath(img_filename)
        img_to_silent(img_filename, silent_filename, mp3_duration)

        video_filename = os.path.join(video_pathname, '{}.mp4'.format(i))
        slient_audio_to_video(silent_filename, mp3_filename, video_filename)
        video_list.append(video_filename)
    merge_vidio(video_list, output_filename)

def slide_to_image(ppt_filename, img_pathname):
    print(ppt_filename, img_pathname)
    if os.path.isdir(img_pathname):
        shutil.rmtree(img_pathname) 
    os.makedirs(img_pathname, exist_ok=True)
    a = os.path.isfile(ppt_filename)
    ppt_filename = os.path.abspath(ppt_filename)
    img_pathname = os.path.abspath(img_pathname)
    utils.save_pptx_as_png(img_pathname, ppt_filename, overwrite_folder=True)

if __name__ == '__main__':
    
    make_video('./data/ciyi/aiqiu_kenqiu', './temp', './output/aiqiu_kenqiu.mp4')

# prs = Presentation('template.pptx')

# for slide in prs.slides:
#     for shape in slide.shapes:
#         if not shape.has_text_frame:
#             continue

#         for para in shape.text_frame.paragraphs:
#             print('--------------', para.text)

#             for run in para.runs:
#                 # if run.text == 'END':
#                 #     run.text = 'ā á ǎ à ō ó ǒ ò ē é ě è ī í ǐ ì ū ú ǔ ù ǖ ǘ ǚ ǜ'
#                 print(run.text)


# prs.save('./out.pptx')
# pptfile = r'C:\Users\liuti\Desktop\digitman\out.pptx'
# png_folder = './output'

# os.makedirs(png_folder, exist_ok=True)
